"""PowerPoint (PPTX) document processor for PII redaction."""

import os
from typing import List, Optional, Callable, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt

from .base import BaseProcessor, BaseRestorer, ProcessingResult
from ..detectors.pii_detector import PIIDetector
from ..detectors.patterns import PIIMatch
from ..storage.mapping_store import MappingStore


class PptxProcessor(BaseProcessor):
    """
    Process PowerPoint (PPTX) files to redact PII.

    Processes all slides, including text in shapes, tables, and notes.
    """

    SUPPORTED_EXTENSIONS = ['.pptx']

    def __init__(
        self,
        detector: PIIDetector,
        mapping_store: MappingStore,
        interactive: bool = True,
        user_callback: Optional[Callable[[PIIMatch], Optional[bool]]] = None,
    ):
        super().__init__(detector, mapping_store, interactive, user_callback)

    def extract_text(self, input_path: str) -> str:
        """Extract all text from a PowerPoint file."""
        prs = Presentation(input_path)
        all_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            all_text.append(f"--- Slide {slide_num} ---")

            # Extract from shapes
            for shape in slide.shapes:
                text = self._extract_shape_text(shape)
                if text:
                    all_text.append(text)

            # Extract from notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text.strip():
                    all_text.append(f"[Notes] {notes_text}")

        return "\n".join(all_text)

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape."""
        texts = []

        # Text frame (most shapes)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        texts.append(run.text)

        # Table
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text)

        return " ".join(texts)

    def process(
        self,
        input_path: str,
        output_path: Optional[str] = None,
    ) -> ProcessingResult:
        """
        Process a PowerPoint file, redacting PII.

        Args:
            input_path: Path to input PowerPoint file
            output_path: Path for output file. If None, creates <input>_redacted.pptx

        Returns:
            ProcessingResult with statistics
        """
        if output_path is None:
            base, ext = os.path.splitext(input_path)
            output_path = f"{base}_redacted{ext}"

        # Load presentation
        prs = Presentation(input_path)

        total_redactions = 0
        redaction_counts = {}
        uncertain_count = 0

        # Process each slide
        for slide in prs.slides:
            # Process shapes
            for shape in slide.shapes:
                stats = self._process_shape(shape)
                total_redactions += stats['redacted']
                uncertain_count += stats['uncertain']
                for pii_type, count in stats['by_type'].items():
                    redaction_counts[pii_type] = redaction_counts.get(pii_type, 0) + count

            # Process notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                stats = self._process_text_frame(slide.notes_slide.notes_text_frame)
                total_redactions += stats['redacted']
                uncertain_count += stats['uncertain']
                for pii_type, count in stats['by_type'].items():
                    redaction_counts[pii_type] = redaction_counts.get(pii_type, 0) + count

        # Save presentation
        prs.save(output_path)

        return ProcessingResult(
            input_path=input_path,
            output_path=output_path,
            total_redactions=total_redactions,
            redactions_by_type=redaction_counts,
            uncertain_count=uncertain_count,
        )

    def _process_shape(self, shape) -> dict:
        """Process a single shape, redacting PII."""
        stats = {'redacted': 0, 'uncertain': 0, 'by_type': {}}

        # Process text frame
        if shape.has_text_frame:
            tf_stats = self._process_text_frame(shape.text_frame)
            stats['redacted'] += tf_stats['redacted']
            stats['uncertain'] += tf_stats['uncertain']
            for pii_type, count in tf_stats['by_type'].items():
                stats['by_type'][pii_type] = stats['by_type'].get(pii_type, 0) + count

        # Process table
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        tf_stats = self._process_text_frame(cell.text_frame)
                        stats['redacted'] += tf_stats['redacted']
                        stats['uncertain'] += tf_stats['uncertain']
                        for pii_type, count in tf_stats['by_type'].items():
                            stats['by_type'][pii_type] = stats['by_type'].get(pii_type, 0) + count

        return stats

    def _process_text_frame(self, text_frame) -> dict:
        """Process a text frame, redacting PII in all runs."""
        stats = {'redacted': 0, 'uncertain': 0, 'by_type': {}}

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue

                # Process run text
                processed_text, run_stats = self.process_text(run.text)

                if run_stats['redacted'] > 0:
                    run.text = processed_text
                    stats['redacted'] += run_stats['redacted']
                    stats['uncertain'] += run_stats['uncertain']
                    for pii_type, count in run_stats['by_type'].items():
                        stats['by_type'][pii_type] = stats['by_type'].get(pii_type, 0) + count

        return stats


class PptxRestorer(BaseRestorer):
    """Restore redacted PowerPoint files using mapping store."""

    def __init__(self, mapping_store: MappingStore):
        super().__init__(mapping_store)

    def restore(
        self,
        input_path: str,
        output_path: Optional[str] = None,
    ) -> Tuple[str, int]:
        """
        Restore a redacted PowerPoint file.

        Args:
            input_path: Path to redacted PowerPoint file
            output_path: Path for restored file

        Returns:
            Tuple of (output_path, restoration_count)
        """
        if output_path is None:
            base, ext = os.path.splitext(input_path)
            if base.endswith('_redacted'):
                base = base[:-9]
            output_path = f"{base}_restored{ext}"

        # Get restorations
        restorations = self.mapping_store.get_all_restorations()
        if not restorations:
            raise ValueError("No mappings available for restoration")

        # Load presentation
        prs = Presentation(input_path)
        restoration_count = 0

        # Process each slide
        for slide in prs.slides:
            # Process shapes
            for shape in slide.shapes:
                count = self._restore_shape(shape, restorations)
                restoration_count += count

            # Process notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                count = self._restore_text_frame(
                    slide.notes_slide.notes_text_frame,
                    restorations
                )
                restoration_count += count

        # Save presentation
        prs.save(output_path)

        return output_path, restoration_count

    def _restore_shape(self, shape, restorations: dict) -> int:
        """Restore placeholders in a shape."""
        count = 0

        if shape.has_text_frame:
            count += self._restore_text_frame(shape.text_frame, restorations)

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        count += self._restore_text_frame(cell.text_frame, restorations)

        return count

    def _restore_text_frame(self, text_frame, restorations: dict) -> int:
        """Restore placeholders in a text frame."""
        count = 0

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                for placeholder, original in restorations.items():
                    if placeholder in text:
                        occurrences = text.count(placeholder)
                        text = text.replace(placeholder, original)
                        count += occurrences

                if text != run.text:
                    run.text = text

        return count
